import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from typing import List, Dict, Any
import os
import asyncio
import markdown
from concurrent.futures import ThreadPoolExecutor
from src_ai.core.logger import logger


def _pdf_extract_page(page_num: int, doc_path: str) -> Dict[str, Any]:
    """One page per thread; each thread opens its own document (PyMuPDF is not thread-safe on one doc)."""
    try:
        doc = fitz.open(doc_path)
        page = doc[page_num]
        text = page.get_text()
        doc.close()
        return {
            "content": text,
            "metadata": {
                "source": os.path.basename(doc_path),
                "page_number": page_num + 1,
                "file_type": "pdf",
            },
        }
    except Exception as e:
        logger.error(f"Error extracting PDF page {page_num} from {doc_path}: {str(e)}")
        return {"content": "", "metadata": {"source": os.path.basename(doc_path), "page_number": page_num + 1, "error": str(e)}}


def _load_pdf_pages_threaded(file_path: str) -> List[Dict[str, Any]]:
    """CPU/IO-bound PDF text extraction parallelized with ThreadPoolExecutor."""
    try:
        doc = fitz.open(file_path)
        num_pages = len(doc)
        doc.close()
    except Exception as e:
        logger.error(f"Failed to open PDF {file_path}: {str(e)}")
        return []

    if num_pages == 0:
        return []

    cpu = os.cpu_count() or 4
    max_workers = min(32, num_pages, max(4, cpu * 2))

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        # map preserves input order → page order stays correct
        results = list(pool.map(_pdf_extract_page, range(num_pages), [file_path] * num_pages))
        return [r for r in results if r["content"].strip()]


class DocumentLoader:
    """Consolidated document loader for all professional file types with Parallel Processing."""

    @staticmethod
    async def load_pdf_async(file_path: str) -> List[Dict[str, Any]]:
        """Loads PDF pages in parallel via ThreadPoolExecutor (non-blocking to asyncio event loop)."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _load_pdf_pages_threaded, file_path)

    @staticmethod
    async def load_table_async(file_path: str) -> List[Dict[str, Any]]:
        """Loads CSV/Excel in an async-friendly way."""
        loop = asyncio.get_event_loop()
        
        def _read_table():
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.csv':
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                
                rows = []
                for index, row in df.iterrows():
                    content = row.to_string()
                    if content.strip():
                        rows.append({
                            "content": content,
                            "metadata": {
                                "source": os.path.basename(file_path),
                                "row_index": index + 1,
                                "file_type": "table"
                            }
                        })
                return rows
            except Exception as e:
                logger.error(f"Error loading table {file_path}: {str(e)}")
                return []

        return await loop.run_in_executor(None, _read_table)

    @staticmethod
    async def load_docx_async(file_path: str) -> List[Dict[str, Any]]:
        """Loads DOCX in an async-friendly way."""
        loop = asyncio.get_event_loop()
        
        def _read_docx():
            try:
                doc = Document(file_path)
                full_text = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        full_text.append(para.text)
                
                if not full_text: return []
                
                return [{
                    "content": "\n".join(full_text),
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "file_type": "docx"
                    }
                }]
            except Exception as e:
                logger.error(f"Error loading DOCX {file_path}: {str(e)}")
                return []

        return await loop.run_in_executor(None, _read_docx)

    @staticmethod
    async def load_pptx_async(file_path: str) -> List[Dict[str, Any]]:
        """Loads PPTX in an async-friendly way."""
        loop = asyncio.get_event_loop()
        
        def _read_pptx():
            try:
                prs = Presentation(file_path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    
                    if slide_text:
                        slides.append({
                            "content": "\n".join(slide_text),
                            "metadata": {
                                "source": os.path.basename(file_path),
                                "slide_number": i + 1,
                                "file_type": "pptx"
                            }
                        })
                return slides
            except Exception as e:
                logger.error(f"Error loading PPTX {file_path}: {str(e)}")
                return []

        return await loop.run_in_executor(None, _read_pptx)

    @staticmethod
    async def load_txt_async(file_path: str) -> List[Dict[str, Any]]:
        """Loads TXT or Markdown in an async-friendly way."""
        loop = asyncio.get_event_loop()
        
        def _read_txt():
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                
                if not content.strip(): return []

                ext = os.path.splitext(file_path)[1].lower()
                file_type = "markdown" if ext == '.md' else "txt"
                
                return [{
                    "content": content,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "file_type": file_type
                    }
                }]
            except Exception as e:
                logger.error(f"Error loading text file {file_path}: {str(e)}")
                return []

        return await loop.run_in_executor(None, _read_txt)

class LoaderFactory:
    """Professional Factory to handle multi-format document ingestion asynchronously."""
    
    @staticmethod
    async def load_async(file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Asynchronously loading {ext} file: {file_path}")
        
        if ext == '.pdf':
            return await DocumentLoader.load_pdf_async(file_path)
        elif ext in ['.csv', '.xlsx', '.xls']:
            return await DocumentLoader.load_table_async(file_path)
        elif ext == '.docx':
            return await DocumentLoader.load_docx_async(file_path)
        elif ext == '.pptx':
            return await DocumentLoader.load_pptx_async(file_path)
        elif ext in ['.txt', '.md']:
            return await DocumentLoader.load_txt_async(file_path)
        else:
            logger.error(f"Unsupported file type: {ext}")
            raise ValueError(f"Unsupported file type: {ext}")
